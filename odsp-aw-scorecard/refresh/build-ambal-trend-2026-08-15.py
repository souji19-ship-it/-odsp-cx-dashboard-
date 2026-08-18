import csv
import json
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = ROOT / "data-pipeline" / "ambal-weekly-trend-2026-08-15.json"
PPTX_PATH = ROOT / "dashboard" / "ODSP-AW-Ambal-Weekly-Trends-Jul5-Aug15-v3.pptx"
PREVIEW_PATH = ROOT / "dashboard" / "ODSP-AW-Ambal-Weekly-Trends-Jul5-Aug15-v3-preview.png"
CSV_PATH = ROOT / "data-pipeline" / "ambal-weekly-trend-2026-08-15.csv"

NAVY = "14263D"
BLUE = "195AA6"
AMBER = "C87800"
GREEN = "087A55"
RED = "C43C3C"
TEXT = "25354A"
MUTED = "6B778A"
GRID = "D6DEE8"
SUBHEAD = "EEF2F7"
ALT = "F8FAFC"
WHITE = "FFFFFF"
LATEST = "FCE4EC"
LATEST_LINE = "C63269"

WEEKS = []
COWORK_LATEST_CONTEXT = {}


def rgb(value):
    return RGBColor.from_string(value)


def format_value(metric, value):
    if value is None:
        return "N/A"
    if "Rate" in metric:
        return f"{value:.1f}%"
    return f"{value:,}"


def delta(metric, values, offset):
    current = values[-1]
    baseline = values[-1 - offset]
    if current is None or baseline is None:
        return "N/A", MUTED
    if "Rate" in metric:
        change = current - baseline
        if abs(change) < 0.05:
            return "≈ 0.0pp", MUTED
        return f"{'▲' if change > 0 else '▼'} {abs(change):.1f}pp", GREEN if change > 0 else RED
    if baseline == 0:
        return "N/A", MUTED
    change = (current - baseline) / baseline * 100
    if abs(change) < 0.05:
        return "≈ 0.0%", MUTED
    return f"{'▲' if change > 0 else '▼'} {abs(change):.1f}%", GREEN if change > 0 else RED


def add_box(slide, x, y, w, h, text="", fill=WHITE, line=GRID, color=TEXT,
            size=8.2, bold=False, align=PP_ALIGN.CENTER, margin=0.025):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.45)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return shape


def add_table(slide, title, y, metrics, values_by_metric):
    x = 0.28
    widths = [0.42, 1.52] + [1.03] * 6 + [0.78, 0.78]
    header_h = 0.34
    row_h = 0.39
    title_h = 0.24
    total_w = sum(widths)

    title_box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(total_w), Inches(title_h)
    )
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = rgb(LATEST_LINE)

    top = y + title_h
    headers = ["Pillar", "Metric"] + WEEKS + ["WoW", "MoM"]
    cursor = x
    for index, (header, width) in enumerate(zip(headers, widths)):
        fill = LATEST if index == 7 else SUBHEAD
        line = LATEST_LINE if index == 7 else GRID
        add_box(slide, cursor, top, width, header_h, header, fill, line, TEXT, 7.5, True)
        cursor += width

    pillar_groups = [
        ("Reach", 0, 3, BLUE),
        ("Adoption", 3, 3 if title == "Copilot Studio" else 2, AMBER),
        ("Reliability", 6 if title == "Copilot Studio" else 5, 1, GREEN),
    ]
    for label, start, count, color in pillar_groups:
        add_box(
            slide, x, top + header_h + start * row_h, widths[0], count * row_h,
            label, color, color, WHITE, 7.2, True
        )

    for row_index, metric in enumerate(metrics):
        row_y = top + header_h + row_index * row_h
        base_fill = ALT if row_index % 2 else WHITE
        cursor = x + widths[0]
        add_box(
            slide, cursor, row_y, widths[1], row_h, metric, base_fill, GRID,
            TEXT, 7.6, True, PP_ALIGN.LEFT, 0.06
        )
        cursor += widths[1]
        values = values_by_metric[metric]
        for week_index, value in enumerate(values):
            fill = LATEST if week_index == len(values) - 1 else base_fill
            line = LATEST_LINE if week_index == len(values) - 1 else GRID
            cell_text = format_value(metric, value)
            cell_size = 7.4
            if (
                title == "Cowork"
                and week_index == len(values) - 1
                and metric in COWORK_LATEST_CONTEXT.get("allUp", {})
            ):
                all_up = COWORK_LATEST_CONTEXT["allUp"][metric]
                share = COWORK_LATEST_CONTEXT["odspShare"][metric]
                cell_text = f"{cell_text}\n({all_up:,} · {share:.1f}%)"
                cell_size = 6.4
            add_box(
                slide, cursor, row_y, widths[2 + week_index], row_h,
                cell_text, fill, line, TEXT, cell_size,
                week_index == len(values) - 1
            )
            cursor += widths[2 + week_index]
        wow, wow_color = delta(metric, values, 1)
        mom, mom_color = delta(metric, values, 4)
        add_box(slide, cursor, row_y, widths[-2], row_h, wow, base_fill, GRID, wow_color, 7.2, True)
        cursor += widths[-2]
        add_box(slide, cursor, row_y, widths[-1], row_h, mom, base_fill, GRID, mom_color, 7.2, True)

    return top + header_h + len(metrics) * row_h


def build_ppt(data):
    global COWORK_LATEST_CONTEXT, WEEKS
    WEEKS = data["weeks"]
    COWORK_LATEST_CONTEXT = data.get("coworkLatestContext", {})
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_box(slide, 0, 0, 13.333, 0.48, "", NAVY, NAVY)
    title = slide.shapes.add_textbox(Inches(0.34), Inches(0.08), Inches(9.8), Inches(0.32))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "ODSP in Agentic Work — Retained Weekly Trends"
    r.font.name = "Aptos Display"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = rgb(WHITE)
    subtitle = slide.shapes.add_textbox(Inches(10.0), Inches(0.10), Inches(3.0), Inches(0.28))
    p = subtitle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "Jul 5–Aug 15, 2026"
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = rgb(WHITE)

    cowork_metrics = [
        "Active Users", "Active Tenants", "Active Agents",
        "Tasks", "Tool Calls", "Tool Success Rate",
    ]
    cs_metrics = [
        "Active Users", "Active Tenants", "Active Agents",
        "Tasks", "Tool Calls", "Knowledge Search", "Tool Success Rate",
    ]
    cowork_bottom = add_table(slide, "Cowork", 0.60, cowork_metrics, data["cowork"])
    cs_bottom = add_table(slide, "Copilot Studio", cowork_bottom + 0.18, cs_metrics, data["copilotStudio"])

    note = slide.shapes.add_textbox(
        Inches(0.30), Inches(cs_bottom + 0.06), Inches(12.7), Inches(0.38)
    )
    tf = note.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (
        "Source: Jul 5–Aug 1 exact values from Ambal's detailed July table. Aug 2–15 Cowork values "
        "use the supplied six-week trailing scorecard (reporting date Aug 15). Aug 2–15 Copilot "
        "Studio uses the supplied platform-team canonical 17-region basis. MoM compares Aug 9–15 "
        "with Jul 12–18."
    )
    r.font.name = "Aptos"
    r.font.size = Pt(7.2)
    r.font.color.rgb = rgb(MUTED)
    prs.save(PPTX_PATH)


def build_preview(data):
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    try:
        regular = ImageFont.truetype("arial.ttf", 19)
        bold = ImageFont.truetype("arialbd.ttf", 20)
        small = ImageFont.truetype("arial.ttf", 14)
    except OSError:
        regular = bold = small = ImageFont.load_default()
    draw.rectangle((0, 0, 1600, 60), fill="#" + NAVY)
    draw.text((38, 16), "ODSP in Agentic Work — Retained Weekly Trends", font=bold, fill="white")

    def render(title, y, metrics, values):
        draw.text((35, y), title, font=bold, fill="#" + LATEST_LINE)
        y += 30
        col_x = [35, 100, 290, 420, 550, 680, 810, 940, 1070, 1200]
        headers = ["Pillar", "Metric"] + data["weeks"] + ["WoW", "MoM"]
        for i, header in enumerate(headers):
            x0 = col_x[i]
            x1 = col_x[i + 1] if i + 1 < len(col_x) else 1555
            fill = "#" + (LATEST if i == 7 else SUBHEAD)
            draw.rectangle((x0, y, x1, y + 30), fill=fill, outline="#" + GRID)
            draw.text((x0 + 5, y + 7), header, font=small, fill="#" + TEXT)
        y += 30
        for row, metric in enumerate(metrics):
            fill = "#" + (ALT if row % 2 else WHITE)
            draw.rectangle((35, y, 1555, y + 38), fill=fill, outline="#" + GRID)
            draw.text((105, y + 8), metric, font=small, fill="#" + TEXT)
            for index, value in enumerate(values[metric]):
                x0 = col_x[index + 2]
                x1 = col_x[index + 3]
                if index == 5:
                    draw.rectangle((x0, y, x1, y + 38), fill="#" + LATEST, outline="#" + LATEST_LINE)
                draw.text((x0 + 5, y + 8), format_value(metric, value), font=small, fill="#" + TEXT)
            wow, _ = delta(metric, values[metric], 1)
            mom, _ = delta(metric, values[metric], 4)
            draw.text((1205, y + 8), wow, font=small, fill="#" + TEXT)
            draw.text((1375, y + 8), mom, font=small, fill="#" + TEXT)
            y += 38
        return y

    cowork_metrics = ["Active Users", "Active Tenants", "Active Agents", "Tasks", "Tool Calls", "Tool Success Rate"]
    cs_metrics = ["Active Users", "Active Tenants", "Active Agents", "Tasks", "Tool Calls", "Knowledge Search", "Tool Success Rate"]
    bottom = render("Cowork", 82, cowork_metrics, data["cowork"])
    render("Copilot Studio", bottom + 25, cs_metrics, data["copilotStudio"])
    image.save(PREVIEW_PATH)


def build_csv(data):
    sections = [
        ("Cowork", data["cowork"]),
        ("Copilot Studio", data["copilotStudio"]),
    ]
    with CSV_PATH.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.writer(handle)
        writer.writerow(["Pillar", "Metric", *data["weeks"], "WoW", "MoM"])
        for pillar, metrics in sections:
            for metric, values in metrics.items():
                wow, _ = delta(metric, values, 1)
                mom, _ = delta(metric, values, 4)
                writer.writerow([
                    pillar,
                    metric,
                    *[format_value(metric, value) for value in values],
                    wow,
                    mom,
                ])


with DATA_PATH.open(encoding="utf-8") as handle:
    retained_data = json.load(handle)

build_ppt(retained_data)
build_preview(retained_data)
build_csv(retained_data)
print(PPTX_PATH)
print(PREVIEW_PATH)
print(CSV_PATH)
